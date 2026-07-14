#!/usr/bin/env python3
"""Create PDF companions for the generated AI-102 courseware."""

from __future__ import annotations

import re
import sys
from io import BytesIO
from pathlib import Path

from pptx import Presentation
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE, MSO_SHAPE_TYPE
from reportlab.lib import colors
from reportlab.lib.enums import TA_CENTER
from reportlab.lib.pagesizes import A4, landscape, letter
from reportlab.lib.styles import ParagraphStyle, getSampleStyleSheet
from reportlab.lib.units import inch
from reportlab.pdfgen import canvas
from reportlab.platypus import Image, PageBreak, Paragraph, SimpleDocTemplate, Spacer, Table, TableStyle

sys.path.insert(0, str(Path(__file__).resolve().parent))
import build_ai102_courseware as src  # noqa: E402


OUT_DIR = src.OUT_DIR
PPT_PDF = OUT_DIR / f"{src.COURSE_TITLE}-v{src.VERSION}.pdf"
LG_PDF = OUT_DIR / f"LG-{src.COURSE_TITLE}.pdf"
LP_PDF = OUT_DIR / f"LP-{src.COURSE_TITLE}.pdf"


def esc(text: str) -> str:
    text = src.clean_text(str(text))
    return text.replace("&", "&amp;").replace("<", "&lt;").replace(">", "&gt;")


def hc(hex_color: str):
    return colors.HexColor("#" + hex_color.lstrip("#"))


def styles():
    ss = getSampleStyleSheet()
    ss.add(ParagraphStyle(name="CoverTitle", parent=ss["Title"], fontName="Helvetica-Bold", fontSize=24, textColor=hc(src.BLUE), alignment=TA_CENTER, leading=30))
    ss.add(ParagraphStyle(name="SmallCenter", parent=ss["Normal"], alignment=TA_CENTER, fontSize=9, textColor=hc(src.GREY)))
    ss.add(ParagraphStyle(name="H1Blue", parent=ss["Heading1"], fontName="Helvetica-Bold", fontSize=15, textColor=hc(src.BLUE), spaceAfter=8))
    ss.add(ParagraphStyle(name="H2Dark", parent=ss["Heading2"], fontName="Helvetica-Bold", fontSize=12, textColor=hc(src.DARK), spaceAfter=6))
    ss.add(ParagraphStyle(name="BodySmall", parent=ss["BodyText"], fontSize=9, leading=12, spaceAfter=5))
    ss.add(ParagraphStyle(name="TocEntry", parent=ss["BodyText"], fontSize=10.5, leading=15, spaceAfter=2, textColor=hc(src.DARK)))
    ss.add(ParagraphStyle(name="CodeBlock", parent=ss["Code"], fontName="Courier", fontSize=8, leading=10, backColor=hc("F3F5F8"), borderPadding=4))
    return ss


def centered_image(path, width, height):
    img = Image(str(path), width=width, height=height, kind="proportional")
    img.hAlign = "CENTER"
    return img


def add_cover(story, kind: str, ss, logos=None):
    logos = logos or []
    story += [
        Spacer(1, 0.78 * inch),
    ]
    try:
        if src.TERTIARY_LOGO.exists():
            story.append(centered_image(src.TERTIARY_LOGO, 1.25 * inch, 1.25 * inch))
            story.append(Spacer(1, 0.08 * inch))
    except Exception:
        pass
    story += [
        Paragraph(src.ORG, ss["SmallCenter"]),
        Paragraph(f"UEN: 201200696W", ss["SmallCenter"]),
        Spacer(1, 0.18 * inch),
        Paragraph(kind.upper(), ss["CoverTitle"]),
        Spacer(1, 0.06 * inch),
        Paragraph("For", ss["SmallCenter"]),
        Spacer(1, 0.12 * inch),
    ]
    try:
        if logos:
            story.append(centered_image(logos[0], 0.82 * inch, 0.82 * inch))
            story.append(Spacer(1, 0.30 * inch))
    except Exception:
        pass
    story += [
        Paragraph(esc(src.COURSE_TITLE), ss["Title"]),
        Spacer(1, 0.08 * inch),
        Paragraph(f"TGS Ref No: {src.COURSE_CODE}", ss["SmallCenter"]),
    ]
    if kind.lower().startswith("learner"):
        story += [
            Spacer(1, 0.08 * inch),
            Paragraph("Conducted by", ss["SmallCenter"]),
            Paragraph(src.ORG, ss["SmallCenter"]),
            Paragraph("UEN: 201200696W", ss["SmallCenter"]),
        ]
    story += [
        Spacer(1, 0.08 * inch),
        Paragraph(f"Version {src.VERSION}", ss["SmallCenter"]),
        PageBreak(),
    ]


def add_pdf_toc(story, title: str, entries, ss):
    story.append(Paragraph("TABLE OF CONTENTS", ss["H1Blue"]))
    rows = [[Paragraph(esc(section), ss["TocEntry"]), Paragraph(esc(page), ss["TocEntry"])] for section, page in entries]
    table = Table(rows, colWidths=[5.35 * inch, 0.55 * inch])
    table.setStyle(
        TableStyle(
            [
                ("VALIGN", (0, 0), (-1, -1), "TOP"),
                ("ALIGN", (1, 0), (1, -1), "RIGHT"),
                ("LEFTPADDING", (0, 0), (-1, -1), 0),
                ("RIGHTPADDING", (0, 0), (-1, -1), 0),
                ("TOPPADDING", (0, 0), (-1, -1), 2),
                ("BOTTOMPADDING", (0, 0), (-1, -1), 2),
            ]
        )
    )
    story.append(table)
    story.append(PageBreak())


def lg_pdf_toc_entries(labs):
    lab_pages = [4, 6, 8, 11, 13, 15, 17, 19, 21, 23]
    entries = [
        ("Table of Contents", "2"),
        ("Course Overview", "3"),
        ("Before You Start", "3"),
    ]
    for lab, page in zip(labs, lab_pages):
        entries.append((f"Lab {lab.num:02d} - {lab.title}", str(page)))
    entries.append(("Cleanup Checklist", "26"))
    return entries


def lp_pdf_toc_entries():
    return [
        ("Table of Contents", "2"),
        ("Course Information", "3"),
        ("Lesson Schedule with Slide Numbers", "3"),
        ("Labs Covered", "4"),
        ("Assessment Notes", "4"),
    ]


def build_lg_pdf(labs, logos):
    ss = styles()
    doc = SimpleDocTemplate(str(LG_PDF), pagesize=A4, rightMargin=54, leftMargin=54, topMargin=54, bottomMargin=54)
    story = []
    add_cover(story, "Learner Guide", ss, logos)
    add_pdf_toc(story, "Learner Guide", lg_pdf_toc_entries(labs), ss)
    story += [
        Paragraph("Course Overview", ss["H1Blue"]),
        Paragraph("This Learner Guide mirrors the lab sequence for the Microsoft Certified Azure AI Engineer Associate AI-102 Training course. Work through the labs in order and complete every validation check before moving on.", ss["BodySmall"]),
        Paragraph(f"<b>Course Code:</b> {esc(src.COURSE_CODE)}", ss["BodySmall"]),
        Paragraph(f"<b>TSC Title:</b> {esc(src.TSC_TITLE)}", ss["BodySmall"]),
        Paragraph(f"<b>TSC Code:</b> {esc(src.TSC_CODE)}", ss["BodySmall"]),
        Paragraph("Before You Start", ss["H1Blue"]),
    ]
    for item in [
        "Sign in to the Azure subscription or Skillable lab environment provided by the trainer.",
        "Create lab notes using the recommended files from the Tools Guide.",
        "Use the same naming convention for resources so cleanup is simple.",
        "Do not paste real confidential, personal, or customer data into lab prompts or documents.",
    ]:
        story.append(Paragraph(f"- {esc(item)}", ss["BodySmall"]))
    for lab in labs:
        story += [PageBreak(), Paragraph(f"Lab {lab.num:02d} - {esc(lab.title)}", ss["H1Blue"])]
        story.append(Paragraph(f"<b>Main Topic:</b> {esc(lab.main_topic)}", ss["BodySmall"]))
        story.append(Paragraph(f"<b>Scenario:</b> {esc(lab.scenario)}", ss["BodySmall"]))
        story.append(Paragraph("Objectives", ss["H2Dark"]))
        for obj in lab.objectives:
            story.append(Paragraph(f"- {esc(obj)}", ss["BodySmall"]))
        story.append(Paragraph("Step-by-Step Lab Guide", ss["H2Dark"]))
        for i, step in enumerate(lab.steps, 1):
            story.append(Paragraph(f"Step {i}: {esc(step.title)}", ss["H2Dark"]))
            for kind, payload in step.blocks:
                if kind in ("p", "bullet"):
                    prefix = "- " if kind == "bullet" else ""
                    story.append(Paragraph(f"{prefix}{esc(payload)}", ss["BodySmall"]))
                elif kind == "code":
                    story.append(Paragraph(esc(payload).replace("\n", "<br/>"), ss["CodeBlock"]))
                elif kind == "table":
                    rows = [[Paragraph(esc(c), ss["BodySmall"]) for c in row] for row in payload]
                    table = Table(rows, repeatRows=1)
                    table.setStyle(base_table_style())
                    story.append(table)
                    story.append(Spacer(1, 4))
            story.append(Paragraph("1. Record the required outputs in your lab notes.<br/>2. Ask the trainer to verify any uncertain configuration or design decision.", ss["BodySmall"]))
        story.append(Paragraph("Validation", ss["H2Dark"]))
        story.append(Paragraph(esc(lab.validation), ss["BodySmall"]))
        story.append(Paragraph("Checkpoint Questions", ss["H2Dark"]))
        for idx, q in enumerate(lab.checkpoints, 1):
            story.append(Paragraph(f"{idx}. {esc(q)}", ss["BodySmall"]))
        story.append(Paragraph("Course Focus", ss["H2Dark"]))
        story.append(Paragraph(esc(lab.focus), ss["BodySmall"]))
        if lab.diagram and lab.diagram.exists():
            story.append(Paragraph("Imported Reference Diagram", ss["H2Dark"]))
            try:
                story.append(Image(str(lab.diagram), width=5.8 * inch, height=2.8 * inch, kind="proportional"))
            except Exception:
                story.append(Paragraph(f"Reference diagram asset: {esc(lab.diagram.name)}", ss["BodySmall"]))
            story.append(Paragraph(f"Imported from source PPT slide {lab.diagram_source_slide}.", ss["SmallCenter"]))
    story += [PageBreak(), Paragraph("Cleanup Checklist", ss["H1Blue"])]
    for item in [
        "Delete training resource groups only when instructed by the trainer.",
        "Remove Foundry projects, model deployments, search indexes, storage accounts, and document processing resources created only for training.",
        "Confirm that no lab keys, endpoint values, or copied sample data remain in shared notes.",
    ]:
        story.append(Paragraph(f"- {esc(item)}", ss["BodySmall"]))
    doc.build(story, onFirstPage=footer, onLaterPages=footer)


def base_table_style():
    return TableStyle(
        [
            ("GRID", (0, 0), (-1, -1), 0.25, hc("DADCE0")),
            ("BACKGROUND", (0, 0), (-1, 0), hc(src.BLUE)),
            ("TEXTCOLOR", (0, 0), (-1, 0), colors.white),
            ("VALIGN", (0, 0), (-1, -1), "TOP"),
            ("LEFTPADDING", (0, 0), (-1, -1), 4),
            ("RIGHTPADDING", (0, 0), (-1, -1), 4),
            ("TOPPADDING", (0, 0), (-1, -1), 3),
            ("BOTTOMPADDING", (0, 0), (-1, -1), 3),
        ]
    )


def build_lp_pdf(labs, slide_index, logos):
    ss = styles()
    doc = SimpleDocTemplate(str(LP_PDF), pagesize=A4, rightMargin=36, leftMargin=36, topMargin=36, bottomMargin=36)
    story = []
    add_cover(story, "Lesson Plan", ss, logos)
    add_pdf_toc(story, "Lesson Plan", lp_pdf_toc_entries(), ss)
    story.append(Paragraph("Course Information", ss["H1Blue"]))
    info = [
        ["Course Title", src.COURSE_TITLE],
        ["Course Code", src.COURSE_CODE],
        ["TSC Title", src.TSC_TITLE],
        ["TSC Code", src.TSC_CODE],
        ["Duration", "2 Days (16 training hours)"],
        ["Daily Schedule", "9:00 AM - 6:00 PM"],
        ["Assessment", "Day 2 final assessment from 4:00 PM"],
    ]
    table = Table([[Paragraph(esc(c), ss["BodySmall"]) for c in row] for row in info], colWidths=[1.35 * inch, 4.8 * inch])
    table.setStyle(base_table_style())
    story += [table, Spacer(1, 10), Paragraph("Lesson Schedule with Slide Numbers", ss["H1Blue"])]
    rows = [["Time", "Topic / Activity", "Duration", "Slides"]]
    rows += [
        ["Day 1 9:00-9:20", "Digital attendance, introduction, ground rules", "20 min", src.slide_range(slide_index, ["cover", "admin_lms_labs"])],
        ["Day 1 9:20-10:20", "Lab 01 - Plan/manage/secure Azure AI solution", "60 min", src.lab_slide_range(slide_index, 1)],
        ["Day 1 10:35-11:35", "Lab 02 - Responsible AI and content safety", "60 min", src.lab_slide_range(slide_index, 2)],
        ["Day 1 11:35-12:35", "Lab 03 - Generative AI, Foundry, RAG", "60 min", src.lab_slide_range(slide_index, 3)],
        ["Day 1 2:00-3:00", "Lab 04 - Agentic AI and Foundry Agent Service", "60 min", src.lab_slide_range(slide_index, 4)],
        ["Day 1 3:00-4:00", "Lab 05 - Vision and video insights", "60 min", src.lab_slide_range(slide_index, 5)],
        ["Day 1 4:15-5:35", "Lab 06 - NLP, speech, translation", "80 min", src.lab_slide_range(slide_index, 6)],
        ["Day 2 9:15-10:15", "Lab 07 - Custom language and Q&A", "60 min", src.lab_slide_range(slide_index, 7)],
        ["Day 2 10:30-11:45", "Lab 08 - AI Search and vector search", "75 min", src.lab_slide_range(slide_index, 8)],
        ["Day 2 11:45-1:00", "Lab 09 - Document Intelligence", "75 min", src.lab_slide_range(slide_index, 9)],
        ["Day 2 2:00-3:30", "Lab 10 - Capstone and course review", "90 min", src.lab_slide_range(slide_index, 10)],
        ["Day 2 4:00-5:15", "Final assessment", "75 min", str(slide_index.get("admin_assessment", "-"))],
        ["Day 2 5:30-6:00", "TRAQOM, support, closing", "30 min", src.slide_range(slide_index, ["closing", "closing_thanks"])],
    ]
    schedule = Table([[Paragraph(esc(c), ss["BodySmall"]) for c in row] for row in rows], colWidths=[1.0 * inch, 3.55 * inch, 0.75 * inch, 0.78 * inch], repeatRows=1)
    schedule.setStyle(base_table_style())
    story.append(schedule)
    story += [PageBreak(), Paragraph("Labs Covered", ss["H1Blue"])]
    lab_rows = [["Lab", "Main Topic", "Title", "PPT Slides"]] + [[f"{lab.num:02d}", lab.main_topic, lab.title, src.lab_slide_range(slide_index, lab.num)] for lab in labs]
    labs_table = Table([[Paragraph(esc(c), ss["BodySmall"]) for c in row] for row in lab_rows], colWidths=[0.42 * inch, 1.35 * inch, 3.55 * inch, 0.80 * inch], repeatRows=1)
    labs_table.setStyle(base_table_style())
    story.append(labs_table)
    story += [Spacer(1, 10), Paragraph("Assessment Notes", ss["H1Blue"])]
    for item in [
        "Ensure AM, PM and Assessment digital attendance are completed.",
        "Assessment is open book using approved materials: slides and Learner Guide.",
        "Learners must complete assessment activities and be assessed as Competent for funding eligibility.",
    ]:
        story.append(Paragraph(f"- {esc(item)}", ss["BodySmall"]))
    doc.build(story, onFirstPage=footer, onLaterPages=footer)


def footer(c, doc):
    c.saveState()
    c.setFont("Helvetica", 7.5)
    c.setFillColor(hc(src.GREY))
    c.drawCentredString(doc.pagesize[0] / 2, 0.35 * inch, f"{src.COPYRIGHT} | Page {doc.page}")
    c.restoreState()


def ppt_shape_text(shape) -> str:
    if hasattr(shape, "text") and shape.text:
        return src.clean_text(shape.text.replace("\n", " | "))
    return ""


def build_ppt_pdf():
    prs = Presentation(str(src.PPTX_OUT))
    page_w, page_h = (13.333 * inch, 7.5 * inch)
    c = canvas.Canvas(str(PPT_PDF), pagesize=(page_w, page_h))
    sx = page_w / prs.slide_width
    sy = page_h / prs.slide_height
    for idx, slide in enumerate(prs.slides, 1):
        c.setFillColor(colors.white)
        c.rect(0, 0, page_w, page_h, fill=1, stroke=0)
        for shape in slide.shapes:
            if shape.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE:
                draw_auto_shape(c, shape, sx, sy, page_h)
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                try:
                    img = Image(BytesIO(shape.image.blob))
                    x = shape.left * sx
                    y = page_h - (shape.top + shape.height) * sy
                    w = shape.width * sx
                    h = shape.height * sy
                    if w > 25 and h > 25:
                        img.drawHeight = h
                        img.drawWidth = w
                        img.wrapOn(c, w, h)
                        img.drawOn(c, x, y)
                except Exception:
                    pass
            text = ppt_shape_text(shape)
            if text:
                draw_text_shape(c, shape, text, sx, sy, page_h)
        c.showPage()
    c.save()


def draw_auto_shape(c, shape, sx, sy, page_h):
    try:
        if shape.auto_shape_type not in (MSO_AUTO_SHAPE_TYPE.RECTANGLE, MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, MSO_AUTO_SHAPE_TYPE.OVAL):
            return
    except Exception:
        return
    try:
        fill_rgb = shape.fill.fore_color.rgb
    except Exception:
        fill_rgb = None
    try:
        line_rgb = shape.line.color.rgb
    except Exception:
        line_rgb = None
    x = shape.left * sx
    y = page_h - (shape.top + shape.height) * sy
    w = shape.width * sx
    h = shape.height * sy
    if not fill_rgb or w <= 0 or h <= 0:
        return
    c.saveState()
    c.setFillColor(colors.HexColor(f"#{fill_rgb}"))
    if line_rgb:
        c.setStrokeColor(colors.HexColor(f"#{line_rgb}"))
        c.setLineWidth(1)
        stroke = 1
    else:
        c.setStrokeColor(colors.HexColor(f"#{fill_rgb}"))
        stroke = 0
    if shape.auto_shape_type == MSO_AUTO_SHAPE_TYPE.OVAL:
        c.ellipse(x, y, x + w, y + h, fill=1, stroke=stroke)
    elif shape.auto_shape_type == MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE:
        c.roundRect(x, y, w, h, min(10, h / 4), fill=1, stroke=stroke)
    else:
        c.rect(x, y, w, h, fill=1, stroke=stroke)
    c.restoreState()


def draw_text_shape(c, shape, text, sx, sy, page_h):
    x = shape.left * sx
    top = shape.top * sy
    w = max(40, shape.width * sx)
    h = max(12, shape.height * sy)
    y = page_h - top - 10
    font_size = first_text_size(shape)
    color = first_text_color(shape)
    bold = first_text_bold(shape)
    c.setFont("Helvetica-Bold" if bold else "Helvetica", font_size)
    c.setFillColor(color)
    max_chars = max(12, int(w / max(font_size * 0.48, 4)))
    lines = []
    for part in text.split(" | "):
        lines.extend(wrap(part, max_chars))
    line_h = font_size + 2
    max_lines = max(1, int(h / line_h) + 1)
    for line in lines[:max_lines]:
        c.drawString(x, y, line)
        y -= line_h


def first_text_size(shape) -> float:
    try:
        for para in shape.text_frame.paragraphs:
            for run in para.runs:
                if run.font.size:
                    return max(6, min(28, run.font.size.pt * 0.72))
    except Exception:
        pass
    return 10


def first_text_color(shape):
    try:
        for para in shape.text_frame.paragraphs:
            for run in para.runs:
                rgb = run.font.color.rgb
                if rgb:
                    return colors.HexColor(f"#{rgb}")
    except Exception:
        pass
    return hc(src.DARK)


def first_text_bold(shape) -> bool:
    try:
        for para in shape.text_frame.paragraphs:
            for run in para.runs:
                if run.font.bold:
                    return True
    except Exception:
        pass
    return False


def wrap(text, chars):
    words = re.split(r"\s+", text)
    lines = []
    current = []
    for word in words:
        if sum(len(x) for x in current) + len(current) + len(word) > chars:
            lines.append(" ".join(current))
            current = [word]
        else:
            current.append(word)
    if current:
        lines.append(" ".join(current))
    return lines


def main():
    OUT_DIR.mkdir(exist_ok=True)
    labs = src.load_labs()
    logos, _ = src.extract_reference_assets(labs)
    builder = src.DeckBuilder(labs, logos)
    slide_index = builder.build()
    build_ppt_pdf()
    build_lg_pdf(labs, logos)
    build_lp_pdf(labs, slide_index, logos)
    print(f"Wrote {PPT_PDF}")
    print(f"Wrote {LG_PDF}")
    print(f"Wrote {LP_PDF}")


if __name__ == "__main__":
    main()
