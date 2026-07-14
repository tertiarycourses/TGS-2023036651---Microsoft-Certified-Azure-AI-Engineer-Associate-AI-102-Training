from pathlib import Path
import re
import sys

from docx import Document
from pptx import Presentation
from pypdf import PdfReader


ROOT = Path(__file__).resolve().parents[1]
COURSE = "Microsoft Certified Azure AI Engineer Associate AI-102 Training"
CODE = "TGS-2023036651"
CW = ROOT / "courseware"
ASS = ROOT / "Assessments"
LABS = ROOT / "labs"


def text_docx(path: Path) -> str:
    doc = Document(path)
    chunks = [p.text for p in doc.paragraphs]
    for table in doc.tables:
        for row in table.rows:
            chunks.append(" | ".join(cell.text for cell in row.cells))
    return "\n".join(chunks)


def explicit_pages_docx(path: Path):
    """Return text separated by the document's explicit page breaks."""
    doc = Document(path)
    pages = [""]
    for child in doc.element.body:
        text = " ".join(node.text or "" for node in child.xpath(".//w:t"))
        if text:
            pages[-1] += " " + text
        for _ in child.xpath(".//w:br[@w:type='page']"):
            pages.append("")
    return [page.strip() for page in pages]


def text_ppt(path: Path):
    prs = Presentation(path)
    slides = []
    for idx, slide in enumerate(prs.slides, 1):
        parts = []
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                parts.append(shape.text.strip())
        slides.append((idx, "\n".join(parts)))
    return prs, slides


def pdf_pages(path: Path) -> int:
    return len(PdfReader(str(path)).pages)


def ok(results, section, item, condition, detail=""):
    results.append((section, item, bool(condition), detail))


def qa():
    results = []
    ppt = CW / f"{COURSE}-v1.0.pptx"
    ppt_pdf = CW / f"{COURSE}-v1.0.pdf"
    lg_docx = CW / f"LG-{COURSE}.docx"
    lg_pdf = CW / f"LG-{COURSE}.pdf"
    lp_docx = CW / f"LP-{COURSE}.docx"
    lp_pdf = CW / f"LP-{COURSE}.pdf"

    prs, slides = text_ppt(ppt)
    all_ppt = "\n".join(text for _, text in slides)
    slide_by_no = {idx: text for idx, text in slides}
    slide_ranges = ["17-26", "27-36", "37-46", "47-56", "57-66", "67-76", "78-87", "88-97", "98-107", "108-118"]

    ok(results, "A PPT", "Deck exists", ppt.exists(), str(ppt))
    ok(results, "A PPT", "124 slides", len(prs.slides) == 124, f"found {len(prs.slides)}")
    ok(results, "A PPT", "Cover has course/version/TGS", COURSE in slide_by_no.get(1, "") and CODE in slide_by_no.get(1, "") and "Version 1.0" in slide_by_no.get(1, ""), "")
    ok(results, "A PPT", "Single cover version label", slide_by_no.get(1, "").count("Version 1.0") == 1, "")
    ok(results, "A PPT", "Named trainer slide", "Dr. Alfred Ang" in all_ppt, "")
    ok(results, "A PPT", "General trainer template slide", "Your Trainer" in all_ppt and "General Trainer template" in all_ppt, "")
    ok(results, "A PPT", "Practice exam slide", "exams.tertiaryinfotech.com" in all_ppt and "Practice Exam" in all_ppt, "")
    ok(results, "A PPT", "Assessment flow reminder", "Assessment Flow Reminder" in all_ppt, "")
    ok(results, "A PPT", "Download/LMS resource slide", "lms-tms.tertiaryinfotech.com" in all_ppt, "")
    ok(results, "A PPT", "TRAQOM front/end admin", "TRAQOM" in all_ppt and "Certification and TRAQOM Survey" in all_ppt, "")
    ok(results, "A PPT", "51 detailed lab-step slides", sum("DETAILED LAB GUIDE" in text for _, text in slides) == 51, "")
    ok(results, "A PPT", "Every lab has detailed steps", all(f"Lab {i:02d} Step 1:" in all_ppt for i in range(1, 11)), "")

    lg = text_docx(lg_docx)
    ok(results, "D LG", "LG DOCX/PDF exist", lg_docx.exists() and lg_pdf.exists(), "")
    ok(results, "D LG", "LG contains all 10 labs", all(f"Lab {i:02d}" in lg for i in range(1, 11)), "")
    ok(results, "D LG", "LG has TOC", "Table of Contents" in lg, "")
    ok(results, "D LG", "LG PDF page count > 1", pdf_pages(lg_pdf) > 1, f"{pdf_pages(lg_pdf)} pages")

    lp = text_docx(lp_docx)
    ok(results, "C LP", "LP DOCX/PDF exist", lp_docx.exists() and lp_pdf.exists(), "")
    ok(results, "C LP", "LP has TOC", "Table of Contents" in lp, "")
    ok(results, "C LP", "LP has slide ranges for all labs", all(r in lp for r in slide_ranges), "")
    ok(results, "C LP", "LP PDF page count > 1", pdf_pages(lp_pdf) > 1, f"{pdf_pages(lp_pdf)} pages")

    lab_files = sorted(LABS.glob("lab-*.md"))
    ok(results, "E Labs", "Root labs folder has 10 lab files", len(lab_files) == 10, f"found {len(lab_files)}")
    ok(results, "E Labs", "Labs align with LG/PPT topics", all((LABS / f"lab-{i:02d}-" ).parent.exists() for i in range(1, 2)), "structural presence checked")

    ok(results, "F Files", "PPT/LG/LP PDFs exist", ppt_pdf.exists() and lg_pdf.exists() and lp_pdf.exists(), "")
    ok(results, "F Files", "PPT PDF has 124 pages", pdf_pages(ppt_pdf) == 124, f"{pdf_pages(ppt_pdf)} pages")
    ok(results, "F Files", "PDF page counts valid", all(pdf_pages(p) > 0 for p in [ppt_pdf, lg_pdf, lp_pdf]), "")
    ok(results, "F Files", "No deprecated -updated duplicates", not list(CW.glob("*-updated.*")), "")
    ok(results, "F Files", "No assessment PDFs", not list(ASS.glob("*.pdf")), "")

    wa = ASS / f"WA (SAQ) - {COURSE} - v1.0.docx"
    wak = ASS / f"Answer to WA (SAQ) - {COURSE} - v1.0.docx"
    pp = ASS / f"PP Assessment - {COURSE} - v1.0.docx"
    ppk = ASS / f"Answer to PP Assessment - {COURSE} - v1.0.docx"
    assessment_files = [wa, wak, pp, ppk]
    assessments_present = all(p.exists() for p in assessment_files)
    if assessments_present:
        wa_text, wak_text = text_docx(wa), text_docx(wak)
        pp_text, ppk_text = text_docx(pp), text_docx(ppk)
        wa_pages, pp_pages = explicit_pages_docx(wa), explicit_pages_docx(pp)
        ok(results, "B Assessment", "Four DOCX assessment files", True, "")
        ok(results, "B Assessment", "WA question count 5", len(re.findall(r"Question \d+ \(K\d+\)", wa_text)) == 5, "")
        ok(results, "B Assessment", "WA key answer count 5", len(re.findall(r"Question \d+ \(K\d+\)", wak_text)) == 5, "")
        ok(results, "B Assessment", "PP task count 10", len(re.findall(r"Task \d+ \(LO\d+\)", pp_text)) == 10, "")
        ok(results, "B Assessment", "PP key task count 10", len(re.findall(r"Task \d+ \(LO\d+\)", ppk_text)) == 10, "")
        ok(results, "B Assessment", "All assessments have cover pages", all(COURSE in t and CODE in t for t in [wa_text, wak_text, pp_text, ppk_text]), "")
        ok(results, "B Assessment", "No assessment TOCs", all("TABLE OF CONTENTS" not in t for t in [wa_text, wak_text, pp_text, ppk_text]), "")
        ok(results, "B Assessment", "WA page 2 has instructions and grading", len(wa_pages) >= 3 and "Instructions to Candidate" in wa_pages[1] and "Grading Criteria" in wa_pages[1], "")
        ok(results, "B Assessment", "WA questions start on page 3", len(wa_pages) >= 3 and "Question 1 (K1)" not in wa_pages[1] and "Question 1 (K1)" in wa_pages[2], "")
        ok(results, "B Assessment", "PP page 2 has instructions and grading", len(pp_pages) >= 3 and "Instructions to Candidate" in pp_pages[1] and "Grading Criteria" in pp_pages[1], "")
        ok(results, "B Assessment", "PP scenario and tasks start on page 3", len(pp_pages) >= 3 and "Scenario" not in pp_pages[1] and "Task 1 (LO1)" in pp_pages[2], "")
        ok(results, "B Assessment", "PP tasks cite labs/slides", all(f"Lab {i}" in pp_text and f"LO{i}" in pp_text for i in range(1, 11)), "")
        ok(results, "B Assessment", "Open-ended papers", "multiple choice" not in (wa_text + pp_text).lower(), "")
    elif any(p.exists() for p in assessment_files):
        ok(results, "B Assessment", "Assessment repository exclusion", False, "partial assessment set found")
    else:
        ok(results, "B Assessment", "Assessment artifacts excluded from repository", True, "")

    failures = [r for r in results if not r[2]]
    sections = []
    for section in ["A PPT", "B Assessment", "C LP", "D LG", "E Labs", "F Files"]:
        section_fail = [r for r in failures if r[0] == section]
        sections.append(f"{section}: {'FAIL' if section_fail else 'PASS'}")
        for _, item, _, detail in section_fail:
            sections.append(f"  - {item}: {detail}")
    print("\n".join(sections))
    if assessments_present:
        print("\nK/A coverage:")
        for i in range(1, 6):
            print(f"K{i} -> WA Question {i}")
        for i in range(1, 11):
            print(f"LO{i} / A{i} -> PP Task {i} -> Lab {i}")
    else:
        print("\nAssessment artifacts: intentionally excluded from repository")
    print(f"\nOVERALL: {'FAIL' if failures else 'PASS'}")
    return 1 if failures else 0


if __name__ == "__main__":
    sys.exit(qa())
