#!/usr/bin/env python3
"""Build WSQ AI-102 courseware from the local lab Markdown and reference PPTX."""

from __future__ import annotations

import os
import re
import shutil
from dataclasses import dataclass, field
from pathlib import Path
from typing import Iterable

from docx import Document
from docx.enum.section import WD_SECTION
from docx.enum.table import WD_TABLE_ALIGNMENT, WD_CELL_VERTICAL_ALIGNMENT
from docx.enum.text import WD_ALIGN_PARAGRAPH, WD_BREAK, WD_TAB_ALIGNMENT, WD_TAB_LEADER
from docx.oxml import OxmlElement
from docx.oxml.ns import qn
from docx.shared import Inches as DocxInches
from docx.shared import Pt, RGBColor
from pptx import Presentation
from pptx.dml.color import RGBColor as PptRGBColor
from pptx.enum.shapes import MSO_SHAPE, MSO_SHAPE_TYPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt as PptPt


ROOT = Path(__file__).resolve().parents[1]
LAB_DIR = ROOT / ".codex-tmp-course" / "labs"
REFERENCE_PPTX = ROOT / "Resources" / "WSQ - Master Trainer Slides - Microsoft Azure AI Engineer Associate (AI-102) - v11.pptx"
OUT_DIR = ROOT / "courseware"
ASSET_DIR = OUT_DIR / "assets"
TERTIARY_LOGO = ASSET_DIR / "tertiary-logo.png"

COURSE_TITLE = "Microsoft Certified Azure AI Engineer Associate AI-102 Training"
COURSE_CODE = "TGS-2023036651"
TSC_TITLE = "Artificial Intelligence Application in Product Development"
TSC_CODE = "ICT-TEM-4034-1.1"
VERSION = "1.0"
VERSION_DATE = "13 Jul 2026"
ORG = "Tertiary Infotech Academy Pte Ltd"
COPYRIGHT = "Copyright 2026 Tertiary Infotech Academy Pte Ltd. All rights reserved."
EXAM_URL = "https://exams.tertiaryinfotech.com"

PPTX_OUT = OUT_DIR / f"{COURSE_TITLE}-v{VERSION}.pptx"
PPTX_UPDATED_OUT = OUT_DIR / f"{COURSE_TITLE}-v{VERSION}-updated.pptx"
LG_DOCX_OUT = OUT_DIR / f"LG-{COURSE_TITLE}.docx"
LG_MD_OUT = OUT_DIR / f"LG-{COURSE_TITLE}.md"
LP_DOCX_OUT = OUT_DIR / f"LP-{COURSE_TITLE}.docx"

BLUE = "1F6FEB"
TEAL = "10B981"
NAVY = "0B1220"
DARK = "111827"
GREY = "5B6372"
LIGHT = "F5F8FC"
LINE = "E2E8F0"
WHITE = "FFFFFF"
AMBER = "F59E0B"
FONT = "Arial"


@dataclass
class LabStep:
    title: str
    blocks: list[tuple[str, object]] = field(default_factory=list)


@dataclass
class Lab:
    num: int
    title: str
    main_topic: str
    objectives: list[str]
    scenario: str
    steps: list[LabStep]
    validation: str
    checkpoints: list[str]
    focus: str
    source_file: Path
    diagram: Path | None = None
    diagram_source_slide: int | None = None


def clean_text(text: str) -> str:
    text = text.replace("\u2013", "-").replace("\u2014", "-").replace("\u2018", "'").replace("\u2019", "'")
    text = text.replace("\u201c", '"').replace("\u201d", '"').replace("\u00a0", " ")
    return re.sub(r"[ \t]+", " ", text).strip()


def slugify(text: str) -> str:
    text = re.sub(r"[^a-zA-Z0-9]+", "-", text).strip("-").lower()
    return text or "item"


def read_lab_index() -> dict[int, str]:
    index = {}
    md = (LAB_DIR / "README.md").read_text(encoding="utf-8")
    for line in md.splitlines():
        m = re.match(r"\|\s*(\d+)\s*\|\s*\[[^\]]+\]\([^)]+\)\s*\|\s*([^|]+)\|", line)
        if m:
            index[int(m.group(1))] = clean_text(m.group(2))
    return index


def split_sections(lines: list[str]) -> dict[str, list[str]]:
    sections: dict[str, list[str]] = {}
    current = ""
    for line in lines:
        if line.startswith("## "):
            current = line[3:].strip()
            sections[current] = []
        elif current:
            sections[current].append(line)
    return sections


def parse_list(lines: Iterable[str]) -> list[str]:
    out = []
    for line in lines:
        s = line.strip()
        m = re.match(r"^[-*]\s+(.*)$", s) or re.match(r"^\d+\.\s+(.*)$", s)
        if m:
            out.append(clean_text(m.group(1)))
    return out


def parse_paragraph(lines: Iterable[str]) -> str:
    parts = []
    for line in lines:
        s = line.strip()
        if not s or s.startswith("#") or s.startswith("|") or s.startswith("```") or s.startswith("- "):
            continue
        if re.match(r"^\d+\.\s+", s):
            continue
        parts.append(clean_text(s))
    return " ".join(parts)


def parse_step_blocks(lines: list[str]) -> list[LabStep]:
    steps: list[LabStep] = []
    current: LabStep | None = None
    in_code = False
    code_lines: list[str] = []
    table_lines: list[str] = []

    def flush_table() -> None:
        nonlocal table_lines
        if current and table_lines:
            rows = []
            for row in table_lines:
                if re.match(r"^\|\s*---", row):
                    continue
                cells = [clean_text(c) for c in row.strip().strip("|").split("|")]
                rows.append(cells)
            if rows:
                current.blocks.append(("table", rows))
        table_lines = []

    for raw in lines:
        line = raw.rstrip()
        if line.startswith("### "):
            flush_table()
            current = LabStep(clean_text(re.sub(r"^\d+\.\s*", "", line[4:].strip())))
            steps.append(current)
            continue
        if current is None:
            continue
        if line.startswith("```"):
            if in_code:
                current.blocks.append(("code", "\n".join(code_lines).strip()))
                code_lines = []
                in_code = False
            else:
                flush_table()
                in_code = True
            continue
        if in_code:
            code_lines.append(line)
            continue
        if line.strip().startswith("|"):
            table_lines.append(line)
            continue
        flush_table()
        s = line.strip()
        if not s:
            continue
        m = re.match(r"^[-*]\s+(.*)$", s) or re.match(r"^\d+\.\s+(.*)$", s)
        if m:
            current.blocks.append(("bullet", clean_text(m.group(1))))
        else:
            current.blocks.append(("p", clean_text(s)))
    flush_table()
    return steps


def load_labs() -> list[Lab]:
    topics = read_lab_index()
    labs = []
    for path in sorted(LAB_DIR.glob("lab-*.md")):
        text = path.read_text(encoding="utf-8").replace("\r\n", "\n")
        lines = text.splitlines()
        title_line = next((l for l in lines if l.startswith("# ")), path.stem)
        m = re.match(r"# Lab\s+(\d+)\s*-\s*(.*)", title_line)
        if not m:
            continue
        num = int(m.group(1))
        title = clean_text(m.group(2))
        sections = split_sections(lines)
        labs.append(
            Lab(
                num=num,
                title=title,
                main_topic=topics.get(num, ""),
                objectives=parse_list(sections.get("Objectives", [])),
                scenario=parse_paragraph(sections.get("Scenario", [])),
                steps=parse_step_blocks(sections.get("Steps", [])),
                validation=parse_paragraph(sections.get("Validation", [])),
                checkpoints=parse_list(sections.get("Checkpoint Questions", [])),
                focus=parse_paragraph(sections.get("Course Focus", [])),
                source_file=path,
            )
        )
    return labs


def ppt_slide_text(slide) -> str:
    texts = []
    for shape in slide.shapes:
        if hasattr(shape, "text") and shape.text:
            texts.append(shape.text.replace("\n", " "))
    return clean_text(" | ".join(texts))


def save_picture(shape, path: Path) -> Path:
    ext = shape.image.ext or "png"
    if path.suffix.lower() != f".{ext.lower()}":
        path = path.with_suffix(f".{ext}")
    path.write_bytes(shape.image.blob)
    return path


def extract_reference_assets(labs: list[Lab]) -> tuple[list[Path], dict[int, str]]:
    ASSET_DIR.mkdir(parents=True, exist_ok=True)
    prs = Presentation(str(REFERENCE_PPTX))
    slide_text = {i: ppt_slide_text(slide) for i, slide in enumerate(prs.slides, 1)}

    logo_paths: list[Path] = []
    for idx, shape in enumerate(prs.slides[0].shapes, 1):
        if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
            out = save_picture(shape, ASSET_DIR / f"reference-cover-logo-{idx}")
            logo_paths.append(out)

    keyword_map = {
        1: ["Azure AI services", "Secure Azure AI services", "Monitor Azure AI services"],
        2: ["Responsible AI", "Content Safety", "content filters"],
        3: ["Retrieval Augmented Generation", "RAG", "Azure OpenAI"],
        4: ["AI Agent Service", "Agent", "Semantic Kernel"],
        5: ["computer vision", "Azure AI Vision", "Analyze images"],
        6: ["Speech", "Translator", "Translate speech", "Analyze text"],
        7: ["question answering", "Language Understanding", "custom text"],
        8: ["Azure AI Search", "knowledge mining", "vector search"],
        9: ["Document Intelligence", "Content Understanding", "Extract text"],
        10: ["AI-102 Study Guide", "Exam Review", "Practice Exam"],
    }

    used_slides = set()
    for lab in labs:
        best = None
        for i, slide in enumerate(prs.slides, 1):
            text = slide_text[i].lower()
            if i in used_slides:
                continue
            if not any(k.lower() in text for k in keyword_map.get(lab.num, [])):
                continue
            pictures = [s for s in slide.shapes if s.shape_type == MSO_SHAPE_TYPE.PICTURE]
            if pictures:
                best = (i, max(pictures, key=lambda s: s.width * s.height))
                break
        if best:
            slide_num, pic = best
            used_slides.add(slide_num)
            lab.diagram = save_picture(pic, ASSET_DIR / f"reference-lab-{lab.num:02d}-slide-{slide_num}")
            lab.diagram_source_slide = slide_num
    return logo_paths, slide_text


def rgb(hex_color: str) -> RGBColor:
    return RGBColor.from_string(hex_color)


def ppt_rgb(hex_color: str) -> PptRGBColor:
    return PptRGBColor.from_string(hex_color)


def set_doc_run_font(run, size=11, color=DARK, bold=False, italic=False, name=FONT):
    run.font.name = name
    run._element.rPr.rFonts.set(qn("w:ascii"), name)
    run._element.rPr.rFonts.set(qn("w:hAnsi"), name)
    run.font.size = Pt(size)
    run.font.color.rgb = rgb(color)
    run.bold = bold
    run.italic = italic


def shade_cell(cell, fill: str):
    tc_pr = cell._tc.get_or_add_tcPr()
    shd = OxmlElement("w:shd")
    shd.set(qn("w:val"), "clear")
    shd.set(qn("w:color"), "auto")
    shd.set(qn("w:fill"), fill)
    tc_pr.append(shd)


def enable_update_fields(doc: Document):
    settings = doc.settings._element
    for existing in settings.findall(qn("w:updateFields")):
        settings.remove(existing)
    update = OxmlElement("w:updateFields")
    update.set(qn("w:val"), "true")
    settings.append(update)


def append_simple_field(paragraph, instruction: str, placeholder: str = "1"):
    fld = OxmlElement("w:fldSimple")
    fld.set(qn("w:instr"), instruction)
    r = OxmlElement("w:r")
    t = OxmlElement("w:t")
    t.text = placeholder
    r.append(t)
    fld.append(r)
    paragraph._p.append(fld)


def add_doc_footer(doc: Document):
    enable_update_fields(doc)
    for section in doc.sections:
        footer = section.footer
        p = footer.paragraphs[0]
        p.alignment = WD_ALIGN_PARAGRAPH.CENTER
        set_doc_run_font(p.add_run("Page "), size=9, color=GREY)
        append_simple_field(p, "PAGE", "1")
        set_doc_run_font(p.add_run(" of "), size=9, color=GREY)
        append_simple_field(p, "NUMPAGES", "1")
        c = footer.add_paragraph()
        c.alignment = WD_ALIGN_PARAGRAPH.CENTER
        set_doc_run_font(c.add_run(COPYRIGHT), size=7.5, color=GREY)


def setup_a4(doc: Document):
    for section in doc.sections:
        section.page_width = DocxInches(8.27)
        section.page_height = DocxInches(11.69)
        section.top_margin = DocxInches(0.75)
        section.bottom_margin = DocxInches(0.75)
        section.left_margin = DocxInches(0.75)
        section.right_margin = DocxInches(0.75)


def setup_doc_styles(doc: Document):
    normal = doc.styles["Normal"]
    normal.font.name = FONT
    normal.font.size = Pt(11)
    for name, size, color in [
        ("Heading 1", 16, BLUE),
        ("Heading 2", 13, DARK),
        ("Heading 3", 11.5, BLUE),
    ]:
        style = doc.styles[name]
        style.font.name = FONT
        style.font.size = Pt(size)
        style.font.color.rgb = rgb(color)
        style.font.bold = True


def add_cover(doc: Document, kind: str, logos: list[Path]):
    spacer = doc.add_paragraph()
    spacer.paragraph_format.space_after = Pt(44)
    if TERTIARY_LOGO.exists():
        p = doc.add_paragraph()
        p.alignment = WD_ALIGN_PARAGRAPH.CENTER
        p.paragraph_format.space_after = Pt(4)
        try:
            p.add_run().add_picture(str(TERTIARY_LOGO), width=DocxInches(1.35))
        except Exception:
            pass
    cover_rows = [
        (ORG, 10.5, True, DARK, 0),
        ("UEN: 201200696W", 8, False, GREY, 12),
        (kind.upper(), 22, True, BLUE, 6),
        ("For", 8.5, False, GREY, 8),
    ]
    for text, size, bold, color, after in cover_rows:
        p = doc.add_paragraph()
        p.alignment = WD_ALIGN_PARAGRAPH.CENTER
        p.paragraph_format.space_after = Pt(after)
        set_doc_run_font(p.add_run(text), size=size, bold=bold, color=color)
    if logos:
        p = doc.add_paragraph()
        p.alignment = WD_ALIGN_PARAGRAPH.CENTER
        p.paragraph_format.space_after = Pt(22)
        try:
            p.add_run().add_picture(str(logos[0]), width=DocxInches(0.85))
        except Exception:
            pass
    for text, size, bold, color, after in [
        (COURSE_TITLE, 17, True, DARK, 8),
        (f"TGS Ref No: {COURSE_CODE}", 8.8, False, DARK, 10),
    ]:
        p = doc.add_paragraph()
        p.alignment = WD_ALIGN_PARAGRAPH.CENTER
        p.paragraph_format.space_after = Pt(after)
        set_doc_run_font(p.add_run(text), size=size, bold=bold, color=color)
    if kind.lower().startswith("learner"):
        for text, size, bold, color, after in [
            ("Conducted by", 8, False, GREY, 0),
            (ORG, 9.5, True, DARK, 0),
            ("UEN: 201200696W", 8, False, GREY, 8),
        ]:
            p = doc.add_paragraph()
            p.alignment = WD_ALIGN_PARAGRAPH.CENTER
            p.paragraph_format.space_after = Pt(after)
            set_doc_run_font(p.add_run(text), size=size, bold=bold, color=color)
    p = doc.add_paragraph()
    p.alignment = WD_ALIGN_PARAGRAPH.CENTER
    p.paragraph_format.space_after = Pt(0)
    set_doc_run_font(p.add_run(f"Version {VERSION}"), size=9.5, bold=True, color=BLUE)
    doc.add_paragraph().add_run().add_break(WD_BREAK.PAGE)


def add_version_table(doc: Document, summary: str):
    h = doc.add_paragraph()
    set_doc_run_font(h.add_run("DOCUMENT VERSION CONTROL RECORD"), size=12, bold=True)
    table = doc.add_table(rows=1, cols=4)
    table.style = "Table Grid"
    table.alignment = WD_TABLE_ALIGNMENT.CENTER
    headers = ["Version Number", "Effective Date of Release", "Summary of Included Changes", "Author"]
    for i, head in enumerate(headers):
        cell = table.rows[0].cells[i]
        shade_cell(cell, BLUE)
        set_doc_run_font(cell.paragraphs[0].add_run(head), size=9.5, color=WHITE, bold=True)
    row = table.add_row().cells
    values = [VERSION, "8 July 2026", summary, ORG]
    for i, val in enumerate(values):
        set_doc_run_font(row[i].paragraphs[0].add_run(val), size=9.5)
    doc.add_paragraph().add_run().add_break(WD_BREAK.PAGE)


def add_toc(doc: Document, entries: list[tuple[str, str]]):
    h = doc.add_paragraph()
    h.paragraph_format.space_after = Pt(14)
    set_doc_run_font(h.add_run("TABLE OF CONTENTS"), size=14, color=BLUE, bold=True)
    for section, page in entries:
        p = doc.add_paragraph()
        p.paragraph_format.space_after = Pt(3)
        p.paragraph_format.tab_stops.add_tab_stop(DocxInches(6.45), WD_TAB_ALIGNMENT.RIGHT, WD_TAB_LEADER.DOTS)
        set_doc_run_font(p.add_run(section), size=10.5, color=DARK)
        set_doc_run_font(p.add_run("\t" + page), size=10.5, color=DARK)
    doc.add_paragraph().add_run().add_break(WD_BREAK.PAGE)


def add_key_value_table(doc: Document, rows: list[tuple[str, str]]):
    table = doc.add_table(rows=0, cols=2)
    table.style = "Table Grid"
    table.alignment = WD_TABLE_ALIGNMENT.CENTER
    for key, val in rows:
        cells = table.add_row().cells
        shade_cell(cells[0], "F1F5FB")
        set_doc_run_font(cells[0].paragraphs[0].add_run(key), size=9.5, bold=True)
        set_doc_run_font(cells[1].paragraphs[0].add_run(val), size=9.5)


def add_doc_table(doc: Document, rows: list[list[str]]):
    if not rows:
        return
    table = doc.add_table(rows=1, cols=len(rows[0]))
    table.style = "Table Grid"
    table.alignment = WD_TABLE_ALIGNMENT.CENTER
    for i, val in enumerate(rows[0]):
        shade_cell(table.rows[0].cells[i], BLUE)
        set_doc_run_font(table.rows[0].cells[i].paragraphs[0].add_run(val), size=9, color=WHITE, bold=True)
    for row_vals in rows[1:]:
        cells = table.add_row().cells
        for i, val in enumerate(row_vals[: len(cells)]):
            set_doc_run_font(cells[i].paragraphs[0].add_run(val), size=9)
            cells[i].vertical_alignment = WD_CELL_VERTICAL_ALIGNMENT.CENTER


def render_inline_doc(paragraph, text: str):
    parts = re.split(r"(\*\*[^*]+\*\*|`[^`]+`)", text)
    for part in parts:
        if not part:
            continue
        if part.startswith("**") and part.endswith("**"):
            set_doc_run_font(paragraph.add_run(part[2:-2]), bold=True)
        elif part.startswith("`") and part.endswith("`"):
            set_doc_run_font(paragraph.add_run(part[1:-1]), name="Consolas", size=9.5, color="C7254E")
        else:
            set_doc_run_font(paragraph.add_run(part))


def learner_guide_toc_entries(labs: list[Lab]) -> list[tuple[str, str]]:
    entries = [
        ("Document Version Control Record", "2"),
        ("Table of Contents", "3"),
        ("Course Overview", "4"),
        ("Before You Start", "4"),
    ]
    page = 5
    for lab in labs:
        entries.append((f"Lab {lab.num:02d} - {lab.title}", str(page)))
        page += 2
    entries.append(("Cleanup Checklist", str(page)))
    return entries


def lesson_plan_toc_entries() -> list[tuple[str, str]]:
    return [
        ("Document Version Control Record", "2"),
        ("Table of Contents", "3"),
        ("Course Information", "4"),
        ("Lesson Schedule with Slide Numbers", "4"),
        ("Labs Covered", "5"),
        ("Assessment Notes", "6"),
    ]


def build_learner_guide(labs: list[Lab], logos: list[Path]):
    doc = Document()
    setup_a4(doc)
    setup_doc_styles(doc)
    add_cover(doc, "Learner Guide", logos)
    add_version_table(doc, "First version aligned to the 10 AI-102 labs, with step-by-step lab guidance and Markdown mirror.")
    add_toc(doc, learner_guide_toc_entries(labs))

    doc.add_paragraph(style="Heading 1").add_run("Course Overview")
    p = doc.add_paragraph()
    render_inline_doc(
        p,
        "This Learner Guide mirrors the lab sequence for the Microsoft Certified Azure AI Engineer Associate AI-102 Training course. Work through the labs in order and complete every validation check before moving on.",
    )
    add_key_value_table(
        doc,
        [
            ("Course Title", COURSE_TITLE),
            ("Course Code", COURSE_CODE),
            ("TSC Title", TSC_TITLE),
            ("TSC Code", TSC_CODE),
            ("Duration", "2 days / 16 training hours"),
            ("Delivery Mode", "Instructor-led training with guided Azure AI design and lab activities"),
            ("Tools", "Azure portal, Azure AI Foundry, Azure Cloud Shell, Visual Studio Code, REST client/Postman, Python or C# SDK examples"),
        ],
    )
    doc.add_paragraph(style="Heading 1").add_run("Before You Start")
    for item in [
        "Sign in to the Azure subscription or Skillable lab environment provided by the trainer.",
        "Create lab notes using the recommended files from the Tools Guide.",
        "Use the same naming convention for resources so cleanup is simple.",
        "Do not paste real confidential, personal, or customer data into lab prompts or documents.",
        "At the end of each lab, complete the validation and checkpoint questions.",
    ]:
        render_inline_doc(doc.add_paragraph(style="List Bullet"), item)

    for lab in labs:
        doc.add_paragraph(style="Heading 1").add_run(f"Lab {lab.num:02d} - {lab.title}")
        add_key_value_table(
            doc,
            [
                ("Main Topic", lab.main_topic),
                ("Source Lab", f"labs/{lab.source_file.name}"),
                ("Reference Diagram", f"Source PPT slide {lab.diagram_source_slide}" if lab.diagram_source_slide else "No imported image asset found for this lab topic"),
            ],
        )
        doc.add_paragraph(style="Heading 2").add_run("Objectives")
        for obj in lab.objectives:
            render_inline_doc(doc.add_paragraph(style="List Bullet"), obj)
        doc.add_paragraph(style="Heading 2").add_run("Scenario")
        render_inline_doc(doc.add_paragraph(), lab.scenario)
        doc.add_paragraph(style="Heading 2").add_run("Step-by-Step Lab Guide")
        for i, step in enumerate(lab.steps, 1):
            doc.add_paragraph(style="Heading 3").add_run(f"Step {i}: {step.title}")
            for kind, payload in step.blocks:
                if kind == "p":
                    render_inline_doc(doc.add_paragraph(), str(payload))
                elif kind == "bullet":
                    render_inline_doc(doc.add_paragraph(style="List Bullet"), str(payload))
                elif kind == "code":
                    p = doc.add_paragraph()
                    shade = OxmlElement("w:shd")
                    shade.set(qn("w:fill"), "F3F5F8")
                    p._p.get_or_add_pPr().append(shade)
                    set_doc_run_font(p.add_run(str(payload)), name="Consolas", size=9)
                elif kind == "table":
                    add_doc_table(doc, payload)  # type: ignore[arg-type]
            render_inline_doc(doc.add_paragraph(style="List Number"), "Record the required outputs in your lab notes.")
            render_inline_doc(doc.add_paragraph(style="List Number"), "Ask the trainer to verify any uncertain configuration or design decision.")
        doc.add_paragraph(style="Heading 2").add_run("Validation")
        render_inline_doc(doc.add_paragraph(), lab.validation)
        doc.add_paragraph(style="Heading 2").add_run("Checkpoint Questions")
        for q in lab.checkpoints:
            render_inline_doc(doc.add_paragraph(style="List Number"), q)
        doc.add_paragraph(style="Heading 2").add_run("Course Focus")
        render_inline_doc(doc.add_paragraph(), lab.focus)
        doc.add_paragraph(style="Heading 2").add_run("Lab Alignment Diagram")
        add_doc_table(doc, [
            ["Lab Input", "Design / Build Activity", "Validation Output"],
            [lab.scenario, " -> ".join(step.title for step in lab.steps), lab.validation],
        ])
        if lab.diagram and lab.diagram.exists():
            doc.add_paragraph(style="Heading 2").add_run("Imported Reference Diagram")
            p = doc.add_paragraph()
            p.alignment = WD_ALIGN_PARAGRAPH.CENTER
            try:
                p.add_run().add_picture(str(lab.diagram), width=DocxInches(5.8))
                cap = doc.add_paragraph()
                cap.alignment = WD_ALIGN_PARAGRAPH.CENTER
                set_doc_run_font(cap.add_run(f"Reference image imported from source PPT slide {lab.diagram_source_slide}."), size=9, color=GREY, italic=True)
            except Exception:
                render_inline_doc(doc.add_paragraph(), f"Reference image imported from source PPT slide {lab.diagram_source_slide}: {lab.diagram.name}")

    doc.add_paragraph(style="Heading 1").add_run("Cleanup Checklist")
    for item in [
        "Delete training resource groups only when instructed by the trainer.",
        "Remove Foundry projects, model deployments, search indexes, storage accounts, and document processing resources created only for training.",
        "Confirm that no lab keys, endpoint values, or copied sample data remain in shared notes.",
    ]:
        render_inline_doc(doc.add_paragraph(style="List Bullet"), item)
    add_doc_footer(doc)
    doc.save(LG_DOCX_OUT)
    LG_MD_OUT.write_text(render_lg_markdown(labs), encoding="utf-8")


def render_lg_markdown(labs: list[Lab]) -> str:
    lines = [
        f"# Learner Guide - {COURSE_TITLE}",
        "",
        f"**Course Code:** {COURSE_CODE}",
        f"**Version:** {VERSION}",
        f"**Conducted by:** {ORG}",
        "",
        "## Document Version Control Record",
        "",
        "| Version Number | Effective Date of Release | Summary of Included Changes | Author |",
        "| --- | --- | --- | --- |",
        f"| {VERSION} | 8 July 2026 | First version aligned to the 10 AI-102 labs, with step-by-step lab guidance and Markdown mirror. | {ORG} |",
        "",
        "## Course Overview",
        "",
        "This Learner Guide mirrors the lab sequence for the Microsoft Certified Azure AI Engineer Associate AI-102 Training course. Work through the labs in order and complete every validation check before moving on.",
        "",
        "## Before You Start",
        "",
        "- Sign in to the Azure subscription or Skillable lab environment provided by the trainer.",
        "- Create lab notes using the recommended files from the Tools Guide.",
        "- Use the same naming convention for resources so cleanup is simple.",
        "- Do not paste real confidential, personal, or customer data into lab prompts or documents.",
        "- At the end of each lab, complete the validation and checkpoint questions.",
        "",
    ]
    for lab in labs:
        lines += [
            f"## Lab {lab.num:02d} - {lab.title}",
            "",
            f"**Main Topic:** {lab.main_topic}",
            f"**Source Lab:** `labs/{lab.source_file.name}`",
            f"**Reference Diagram:** {'Source PPT slide ' + str(lab.diagram_source_slide) if lab.diagram_source_slide else 'No imported image asset found for this lab topic'}",
            "",
            "### Objectives",
            "",
        ]
        lines += [f"- {x}" for x in lab.objectives] + ["", "### Scenario", "", lab.scenario, "", "### Step-by-Step Lab Guide", ""]
        for i, step in enumerate(lab.steps, 1):
            lines += [f"#### Step {i}: {step.title}", ""]
            for kind, payload in step.blocks:
                if kind == "p":
                    lines += [str(payload), ""]
                elif kind == "bullet":
                    lines += [f"- {payload}"]
                elif kind == "code":
                    lang = "bash" if "az " in str(payload) else "text"
                    lines += ["", f"```{lang}", str(payload), "```", ""]
                elif kind == "table":
                    rows = payload  # type: ignore[assignment]
                    if rows:
                        lines.append("| " + " | ".join(rows[0]) + " |")
                        lines.append("| " + " | ".join("---" for _ in rows[0]) + " |")
                        for row in rows[1:]:
                            lines.append("| " + " | ".join(row) + " |")
                        lines.append("")
            lines += ["1. Record the required outputs in your lab notes.", "2. Ask the trainer to verify any uncertain configuration or design decision.", ""]
        lines += ["### Validation", "", lab.validation, "", "### Checkpoint Questions", ""]
        lines += [f"{i}. {q}" for i, q in enumerate(lab.checkpoints, 1)]
        lines += ["", "### Course Focus", "", lab.focus, ""]
        lines += [
            "### Lab Alignment Diagram",
            "",
            "| Lab Input | Design / Build Activity | Validation Output |",
            "| --- | --- | --- |",
            f"| {lab.scenario} | {' -> '.join(step.title for step in lab.steps)} | {lab.validation} |",
            "",
        ]
        if lab.diagram:
            rel = lab.diagram.relative_to(OUT_DIR).as_posix()
            lines += ["### Imported Reference Diagram", "", f"![Reference diagram from source PPT slide {lab.diagram_source_slide}]({rel})", ""]
    lines += [
        "## Cleanup Checklist",
        "",
        "- Delete training resource groups only when instructed by the trainer.",
        "- Remove Foundry projects, model deployments, search indexes, storage accounts, and document processing resources created only for training.",
        "- Confirm that no lab keys, endpoint values, or copied sample data remain in shared notes.",
        "",
    ]
    return "\n".join(lines)


class DeckBuilder:
    def __init__(self, labs: list[Lab], logos: list[Path]):
        self.labs = labs
        self.logos = logos
        self.prs = Presentation()
        self.prs.slide_width = Inches(13.333)
        self.prs.slide_height = Inches(7.5)
        self.blank = self.prs.slide_layouts[6]
        self.slide_index: dict[str, int] = {}

    def add(self, key: str | None = None):
        slide = self.prs.slides.add_slide(self.blank)
        if key:
            self.slide_index[key] = len(self.prs.slides)
        self.bg(slide)
        return slide

    def bg(self, slide, fill=WHITE):
        shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, self.prs.slide_width, self.prs.slide_height)
        shape.fill.solid()
        shape.fill.fore_color.rgb = ppt_rgb(fill)
        shape.line.fill.background()

    def text(self, slide, x, y, w, h, text, size=18, color=DARK, bold=False, align=PP_ALIGN.LEFT, valign=MSO_ANCHOR.TOP):
        box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
        tf = box.text_frame
        tf.clear()
        tf.word_wrap = True
        tf.margin_left = PptPt(0)
        tf.margin_right = PptPt(0)
        tf.vertical_anchor = valign
        p = tf.paragraphs[0]
        p.alignment = align
        r = p.add_run()
        r.text = clean_text(text)
        r.font.name = FONT
        r.font.size = PptPt(size)
        r.font.bold = bold
        r.font.color.rgb = ppt_rgb(color)
        return box

    def bullets(self, slide, x, y, w, h, items, size=17, color=DARK):
        box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
        tf = box.text_frame
        tf.clear()
        tf.word_wrap = True
        for idx, item in enumerate(items):
            p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
            p.text = clean_text(item)
            p.level = 0
            p.font.name = FONT
            p.font.size = PptPt(size)
            p.font.color.rgb = ppt_rgb(color)
            p.space_after = PptPt(8)
        return box

    def lab_table(self, slide, x, y, w, h, rows):
        if not rows:
            return
        col_count = max(len(row) for row in rows)
        normalized = [list(row) + [""] * (col_count - len(row)) for row in rows]
        table = slide.shapes.add_table(
            len(normalized),
            col_count,
            Inches(x),
            Inches(y),
            Inches(w),
            Inches(h),
        ).table
        for row_idx, row in enumerate(normalized):
            for col_idx, value in enumerate(row):
                cell = table.cell(row_idx, col_idx)
                cell.text = clean_text(str(value))
                cell.margin_left = Inches(0.08)
                cell.margin_right = Inches(0.08)
                cell.margin_top = Inches(0.05)
                cell.margin_bottom = Inches(0.05)
                cell.fill.solid()
                cell.fill.fore_color.rgb = ppt_rgb(BLUE if row_idx == 0 else WHITE)
                for paragraph in cell.text_frame.paragraphs:
                    paragraph.font.name = FONT
                    paragraph.font.size = PptPt(11 if len(normalized) <= 7 else 10)
                    paragraph.font.bold = row_idx == 0
                    paragraph.font.color.rgb = ppt_rgb(WHITE if row_idx == 0 else DARK)

    def code_block(self, slide, x, y, w, h, code):
        self.rounded(slide, x, y, w, h, "F3F6FA", "DCE5F0")
        box = self.text(slide, x + 0.22, y + 0.16, w - 0.44, h - 0.28, code, 12, DARK)
        for paragraph in box.text_frame.paragraphs:
            for run in paragraph.runs:
                run.font.name = "Consolas"
        return box

    def title_slide(self):
        slide = self.add("cover")
        self.rect(slide, 0, 0, 13.333, 0.18, BLUE, BLUE)
        self.rect(slide, 0, 7.32, 13.333, 0.18, TEAL, TEAL)
        if TERTIARY_LOGO.exists():
            try:
                slide.shapes.add_picture(str(TERTIARY_LOGO), Inches(0.63), Inches(0.55), height=Inches(1.1))
            except Exception:
                pass
        if self.logos:
            try:
                slide.shapes.add_picture(str(self.logos[0]), Inches(11.35), Inches(0.58), height=Inches(1.0))
            except Exception:
                pass
        self.text(slide, 0.7, 2.35, 11.8, 0.35, "LEARNER GUIDE · COURSE SLIDES", 16, BLUE, True)
        self.text(slide, 0.7, 2.9, 10.5, 1.45, COURSE_TITLE, 38, DARK, True)
        self.rect(slide, 0.72, 4.55, 2.6, 0.07, TEAL, TEAL)
        self.text(slide, 0.7, 4.95, 11.8, 0.32, f"WSQ Course Code: {COURSE_CODE}", 15, GREY)
        self.text(slide, 0.7, 5.35, 11.8, 0.32, f"WSQ TSC: {TSC_TITLE} ({TSC_CODE})", 14, GREY)
        self.text(slide, 0.7, 5.75, 11.8, 0.32, f"Conducted by {ORG}  ·  UEN 201200696W", 14, GREY)
        self.text(slide, 0.7, 6.22, 4.2, 0.35, f"Version v{VERSION} · Learner Guide Slides", 16, DARK)
        self.rect(slide, 9.0, 5.62, 2.85, 0.54, "EAF3FF", "BBD7F2")
        self.text(slide, 9.2, 5.73, 2.5, 0.25, "Version v1.0 | 08 Jul 2026", 15, DARK)

    def rect(self, slide, x, y, w, h, fill, line=LINE):
        shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
        shape.fill.solid()
        shape.fill.fore_color.rgb = ppt_rgb(fill)
        if line:
            shape.line.color.rgb = ppt_rgb(line)
            shape.line.width = PptPt(0.8)
        else:
            shape.line.fill.background()
        return shape

    def header(self, slide, title, kicker=""):
        self.rect(slide, 0.55, 0.58, 0.12, 0.62, BLUE, BLUE)
        if kicker:
            self.text(slide, 0.82, 0.48, 11.6, 0.28, kicker.upper(), 12, BLUE, True)
            self.text(slide, 0.82, 0.78, 11.7, 0.62, title, 27, DARK, True)
        else:
            self.text(slide, 0.82, 0.62, 11.7, 0.62, title, 28, DARK, True)
        self.rect(slide, 0.82, 1.48, 11.8, 0.02, LINE, LINE)

    def footer(self, slide):
        n = len(self.prs.slides)
        self.text(slide, 0.55, 7.08, 5.2, 0.25, f"{COURSE_TITLE} | {COURSE_CODE}", 8, GREY)
        self.text(slide, 5.1, 7.08, 4.5, 0.25, COPYRIGHT, 7.5, GREY, False, PP_ALIGN.CENTER)
        self.text(slide, 12.0, 7.08, 0.8, 0.25, str(n), 8.5, GREY, True, PP_ALIGN.RIGHT)

    def section(self, title, kicker, key=None, subtitle=""):
        slide = self.add(key)
        self.rect(slide, 0, 0, 0.28, 7.5, BLUE, BLUE)
        self.text(slide, 1.1, 2.35, 11, 0.45, kicker.upper(), 17, BLUE, True)
        self.text(slide, 1.1, 2.9, 11, 1.1, title, 39, DARK, True)
        if subtitle:
            self.text(slide, 1.12, 4.25, 11, 0.55, subtitle, 16, GREY)
        self.footer(slide)

    def content(self, title, items, key=None, kicker="", size=17):
        slide = self.add(key)
        self.header(slide, title, kicker)
        self.bullets(slide, 0.9, 1.9, 11.5, 4.8, items, size=size)
        self.footer(slide)

    def two_col(self, title, left_title, left_items, right_title, right_items, key=None, kicker=""):
        slide = self.add(key)
        self.header(slide, title, kicker)
        self.rect(slide, 0.85, 1.9, 5.65, 4.75, LIGHT)
        self.rect(slide, 6.85, 1.9, 5.65, 4.75, LIGHT)
        self.text(slide, 1.1, 2.15, 5.1, 0.35, left_title, 16, BLUE, True)
        self.text(slide, 7.1, 2.15, 5.1, 0.35, right_title, 16, TEAL, True)
        self.bullets(slide, 1.1, 2.75, 5.0, 3.5, left_items, 14)
        self.bullets(slide, 7.1, 2.75, 5.0, 3.5, right_items, 14)
        self.footer(slide)

    def about_trainer(self, key=None, blank=False):
        slide = self.add(key)
        self.header(slide, "About the Trainer", "YOUR TRAINER")
        self.rect(slide, 0.9, 1.95, 3.5, 3.9, DARK, None)
        if self.logos:
            try:
                slide.shapes.add_picture(str(self.logos[0]), Inches(1.75), Inches(2.35), height=Inches(1.15))
            except Exception:
                pass
        self.text(slide, 1.15, 3.95, 3.0, 0.45, "[ Trainer Name ]" if blank else "Trainer / Facilitator", 18, WHITE, True, PP_ALIGN.CENTER)
        self.text(slide, 1.15, 4.45, 3.0, 0.45, "[ Title / Role ]" if blank else "Azure AI and Cloud Trainer", 13, BLUE, True, PP_ALIGN.CENTER)
        rows = [
            ("Qualifications", "[ Add certifications and qualifications ]" if blank else "Microsoft Azure AI, cloud and application integration expertise."),
            ("Expertise", "[ Add domain expertise ]" if blank else "Azure AI Services, Azure AI Foundry, Azure OpenAI, AI Search, Document Intelligence and responsible AI."),
            ("Experience", "[ Add project / industry experience ]" if blank else "Instructor-led enterprise training, hands-on labs and assessment facilitation."),
            ("Contact / Profile", "[ Add contact or profile link ]" if blank else "Shared by trainer during class."),
        ]
        y = 2.0
        for label, body in rows:
            self.rect(slide, 4.95, y, 7.35, 0.82, LIGHT, LINE)
            self.text(slide, 5.15, y + 0.12, 2.1, 0.26, label, 13, BLUE, True)
            self.text(slide, 7.05, y + 0.12, 4.95, 0.45, body, 12, GREY)
            y += 0.98
        self.footer(slide)

    def flow_diagram(self, slide, x, y, w, labels, color=BLUE):
        gap = 0.28
        box_w = (w - gap * (len(labels) - 1)) / len(labels)
        for idx, label in enumerate(labels):
            bx = x + idx * (box_w + gap)
            self.rect(slide, bx, y, box_w, 0.8, LIGHT, LINE)
            self.text(slide, bx + 0.12, y + 0.18, box_w - 0.24, 0.35, label, 11, DARK, True, PP_ALIGN.CENTER)
            if idx < len(labels) - 1:
                self.text(slide, bx + box_w + 0.03, y + 0.22, 0.22, 0.3, ">", 14, color, True, PP_ALIGN.CENTER)

    def diagram_slide(self, lab: Lab, key: str):
        slide = self.add(key)
        title = "Imported Reference Diagram" if lab.diagram and lab.diagram.exists() else "Lab Alignment Diagram"
        self.header(slide, f"Lab {lab.num:02d}: {title}", lab.main_topic)
        if lab.diagram and lab.diagram.exists():
            try:
                slide.shapes.add_picture(str(lab.diagram), Inches(0.95), Inches(1.85), width=Inches(7.2), height=Inches(4.4))
            except Exception:
                self.text(slide, 1.0, 2.4, 7.2, 0.5, lab.diagram.name, 16, GREY)
            self.text(slide, 8.55, 2.0, 3.8, 0.45, "How to use it", 18, BLUE, True)
            self.bullets(
                slide,
                8.55,
                2.6,
                3.8,
                2.9,
                [
                    f"Connect the diagram to the lab scenario: {lab.scenario[:110]}...",
                    "Use it as a reference while completing the design steps.",
                    f"Imported from source PPT slide {lab.diagram_source_slide}.",
                ],
                size=13,
            )
        else:
            self.text(slide, 0.95, 1.95, 11.4, 0.45, "Lab-aligned design flow", 18, BLUE, True)
            labels = ["Scenario", "Plan", "Configure / Design", "Validate", "Review"]
            self.flow_diagram(slide, 0.95, 2.65, 11.35, labels)
            self.text(slide, 0.95, 3.75, 2.0, 0.35, "Scenario", 14, BLUE, True)
            self.text(slide, 2.75, 3.75, 9.45, 0.75, lab.scenario, 13, DARK)
            self.text(slide, 0.95, 4.8, 2.0, 0.35, "Lab steps", 14, BLUE, True)
            self.bullets(slide, 2.75, 4.75, 9.45, 1.25, [step.title for step in lab.steps], size=11)
            self.text(slide, 0.95, 6.2, 2.0, 0.35, "Validation", 14, BLUE, True)
            self.text(slide, 2.75, 6.2, 9.45, 0.45, lab.validation, 12, GREY)
        self.footer(slide)

    def lab_slides(self, lab: Lab):
        self.content(
            f"Lab {lab.num:02d}: {lab.title}",
            [f"Main topic: {lab.main_topic}", f"Scenario: {lab.scenario}", f"Validation: {lab.validation}"],
            key=f"lab{lab.num:02d}_overview",
            kicker="LAB OVERVIEW",
            size=15,
        )
        self.content(
            f"Lab {lab.num:02d}: Objectives",
            lab.objectives,
            key=f"lab{lab.num:02d}_objectives",
            kicker="WHAT YOU WILL DO",
            size=17,
        )
        self.diagram_slide(lab, f"lab{lab.num:02d}_diagram")
        step_lines = [f"{i}. {step.title}" for i, step in enumerate(lab.steps, 1)]
        self.content(
            f"Lab {lab.num:02d}: Step Sequence",
            step_lines,
            key=f"lab{lab.num:02d}_steps",
            kicker="FOLLOW IN ORDER",
            size=15,
        )
        for step_num, step in enumerate(lab.steps, 1):
            self.lab_step_detail(lab, step_num, step)
        self.content(
            f"Lab {lab.num:02d}: Validation and Checkpoints",
            [lab.validation] + lab.checkpoints,
            key=f"lab{lab.num:02d}_validation",
            kicker="CHECK BEFORE MOVING ON",
            size=14,
        )

    def lab_step_detail(self, lab: Lab, step_num: int, step: LabStep):
        slide = self.add(f"lab{lab.num:02d}_step{step_num:02d}")
        self.header(
            slide,
            f"Lab {lab.num:02d} Step {step_num}: {step.title}",
            f"DETAILED LAB GUIDE | STEP {step_num} OF {len(lab.steps)}",
        )
        self.text(slide, 0.92, 1.68, 11.5, 0.32, lab.title, 11, GREY)
        y = 2.08
        for kind, payload in step.blocks:
            if kind == "p":
                is_label = str(payload).strip().endswith(":") and len(str(payload)) < 60
                self.text(slide, 0.95, y, 11.35, 0.36, str(payload), 14, BLUE if is_label else DARK, is_label)
                y += 0.42
            elif kind == "bullet":
                self.text(slide, 0.99, y - 0.01, 0.20, 0.32, "-", 17, BLUE, True, PP_ALIGN.CENTER)
                self.text(slide, 1.30, y, 10.85, 0.40, str(payload), 13.5, DARK)
                y += 0.47
            elif kind == "code":
                line_count = max(1, str(payload).count("\n") + 1)
                block_h = min(2.70, max(0.82, 0.33 * line_count + 0.34))
                self.code_block(slide, 0.95, y + 0.02, 11.35, block_h, str(payload))
                y += block_h + 0.18
            elif kind == "table":
                row_count = max(1, len(payload))
                table_h = min(3.75, max(1.15, row_count * 0.43))
                self.lab_table(slide, 0.95, y + 0.02, 11.35, table_h, payload)
                y += table_h + 0.18
        self.rounded(slide, 0.95, 6.42, 11.35, 0.42, "EAF3FF", "BBD7F2")
        self.text(
            slide,
            1.15,
            6.53,
            10.95,
            0.20,
            "Record the required output in your lab notes and verify it before continuing.",
            11,
            BLUE,
            True,
        )
        self.footer(slide)

    def rounded(self, slide, x, y, w, h, fill=LIGHT, line=LINE):
        shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
        shape.fill.solid()
        shape.fill.fore_color.rgb = ppt_rgb(fill)
        if line:
            shape.line.color.rgb = ppt_rgb(line)
            shape.line.width = PptPt(0.8)
        else:
            shape.line.fill.background()
        return shape

    def circle(self, slide, x, y, d, fill=BLUE, line=None):
        shape = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(x), Inches(y), Inches(d), Inches(d))
        shape.fill.solid()
        shape.fill.fore_color.rgb = ppt_rgb(fill)
        if line:
            shape.line.color.rgb = ppt_rgb(line)
        else:
            shape.line.fill.background()
        return shape

    def title_slide(self):
        slide = self.add("cover")
        self.rect(slide, 0, 0, 13.333, 0.18, BLUE, BLUE)
        self.rect(slide, 0, 7.32, 13.333, 0.18, TEAL, TEAL)
        if TERTIARY_LOGO.exists():
            try:
                slide.shapes.add_picture(str(TERTIARY_LOGO), Inches(0.65), Inches(0.55), height=Inches(1.05))
            except Exception:
                pass
        if self.logos:
            try:
                slide.shapes.add_picture(str(self.logos[0]), Inches(11.25), Inches(0.58), height=Inches(1.0))
            except Exception:
                pass
        self.text(slide, 0.72, 2.22, 11.8, 0.35, "COURSE SLIDES", 16, BLUE, True)
        self.text(slide, 0.72, 2.82, 10.5, 1.45, COURSE_TITLE, 38, DARK, True)
        self.rect(slide, 0.74, 4.52, 2.65, 0.07, TEAL, TEAL)
        self.text(slide, 0.72, 4.90, 11.8, 0.32, f"TGS Ref No: {COURSE_CODE}", 15, GREY)
        self.text(slide, 0.72, 5.30, 11.8, 0.32, f"TSC: {TSC_TITLE} ({TSC_CODE})", 14, GREY)
        self.text(slide, 0.72, 5.70, 11.8, 0.32, f"Conducted by {ORG} | UEN 201200696W", 14, GREY)
        self.rounded(slide, 9.08, 5.62, 2.85, 0.54, "EAF3FF", "BBD7F2")
        self.text(slide, 9.28, 5.75, 2.45, 0.23, f"Version {VERSION} | {VERSION_DATE}", 13, DARK)
        self.text(slide, 0.72, 7.05, 5.2, 0.22, f"{COURSE_TITLE} | {COURSE_CODE}", 7.5, GREY)
        self.text(slide, 5.35, 7.05, 3.8, 0.22, COPYRIGHT, 6.5, GREY, False, PP_ALIGN.CENTER)

    def about_trainer(self, key=None, blank=False):
        slide = self.add(key)
        if blank:
            self.rect(slide, 0.64, 0.64, 0.10, 0.66, GREY, GREY)
            self.text(slide, 0.94, 0.68, 9.5, 0.24, "YOUR TRAINER - GENERAL", 12, GREY, True)
            self.text(slide, 0.94, 1.05, 9.5, 0.48, "About the Trainer", 28, DARK, True)
            self.rect(slide, 0.92, 1.64, 11.95, 0.02, LINE, LINE)
            self.rect(slide, 1.0, 2.02, 3.65, 4.62, LIGHT, None)
            self.rect(slide, 1.0, 2.02, 3.65, 0.10, GREY, GREY)
            self.circle(slide, 1.88, 2.56, 1.72, GREY)
            self.text(slide, 2.48, 3.05, 0.52, 0.5, "?", 28, WHITE, True, PP_ALIGN.CENTER)
            self.text(slide, 1.35, 4.72, 2.95, 0.35, "Your Trainer", 19, DARK, True, PP_ALIGN.CENTER)
            self.text(slide, 1.52, 5.28, 2.6, 0.55, "General Trainer template -\nto be completed by the trainer", 11, GREY, True, PP_ALIGN.CENTER)
            fields = [
                ("NAME", BLUE),
                ("TITLE / DESIGNATION", TEAL),
                ("QUALIFICATIONS", "7C3AED"),
                ("AREAS OF EXPERTISE", AMBER),
                ("TRAINING & INDUSTRY EXPERIENCE", BLUE),
                ("CONTACT", TEAL),
            ]
            y = 2.02
            for label, color in fields:
                self.rect(slide, 5.05, y, 0.08, 0.62, color, color)
                self.rect(slide, 5.13, y, 7.55, 0.62, LIGHT, None)
                self.text(slide, 5.38, y + 0.16, 4.5, 0.22, label, 10, color, True)
                self.rect(slide, 5.38, y + 0.45, 4.45, 0.01, "D6E1EF", "D6E1EF")
                y += 0.8
        else:
            self.header(slide, "About the Trainer", "YOUR TRAINER")
            self.rect(slide, 1.03, 2.1, 3.65, 4.6, LIGHT, LINE)
            self.rect(slide, 1.03, 2.1, 3.65, 0.10, BLUE, BLUE)
            self.circle(slide, 1.92, 2.64, 1.72, BLUE)
            self.text(slide, 2.47, 3.18, 0.62, 0.38, "AA", 27, WHITE, True, PP_ALIGN.CENTER)
            self.text(slide, 1.38, 4.78, 2.95, 0.34, "Dr. Alfred Ang", 18, DARK, True, PP_ALIGN.CENTER)
            self.text(slide, 1.48, 5.36, 2.75, 0.52, "Principal Trainer\nTertiary Infotech Academy Pte. Ltd.", 11, GREY, True, PP_ALIGN.CENTER)
            rows = [
                ("ROLE", "Principal Trainer, Tertiary Infotech Academy Pte. Ltd.", BLUE),
                ("QUALIFICATIONS", "PhD - specialises in AI, automation and software engineering.", TEAL),
                ("DELIVERS", "WSQ courses on AI agents, automation (n8n) and app development.", "7C3AED"),
                ("FOUNDER", "Founder and lead instructor at Tertiary Infotech / Tertiary Courses.", AMBER),
            ]
            y = 2.1
            for label, body, color in rows:
                self.rect(slide, 5.05, y, 0.08, 1.0, color, color)
                self.rect(slide, 5.13, y, 7.6, 1.0, LIGHT, LINE)
                self.text(slide, 5.36, y + 0.22, 6.9, 0.20, label, 10, color, True)
                self.text(slide, 5.36, y + 0.50, 6.9, 0.34, body, 13, DARK)
                y += 1.22
        self.footer(slide)

    def ground_rules_slide(self, key=None):
        slide = self.add(key)
        self.header(slide, "Ground Rules", "HOUSEKEEPING")
        items = [
            ("1", "Set your mobile phone to silent mode.", BLUE),
            ("2", "Participate actively - no question is too small.", TEAL),
            ("3", "Mutual respect: agree to disagree.", "7C3AED"),
            ("4", "One conversation at a time.", AMBER),
            ("5", "Be punctual; return from breaks on time.", "E11D48"),
            ("6", "75% attendance is required.", "0EA5E9"),
        ]
        positions = [(0.7, 1.86), (6.88, 1.86), (0.7, 3.65), (6.88, 3.65), (0.7, 5.44), (6.88, 5.44)]
        for (num, body, color), (x, y) in zip(items, positions):
            self.rect(slide, x, y, 0.08, 1.50, color, color)
            self.rounded(slide, x + 0.08, y, 5.80, 1.50, "F3F6FA", "DCE5F0")
            self.circle(slide, x + 0.32, y + 0.45, 0.62, color)
            self.text(slide, x + 0.52, y + 0.58, 0.22, 0.25, num, 13, WHITE, True, PP_ALIGN.CENTER)
            self.text(slide, x + 1.15, y + 0.58, 4.25, 0.38, body, 15, DARK)
        self.footer(slide)

    def assessment_flow_reminder(self, key=None):
        slide = self.add(key)
        self.header(slide, "Assessment Flow Reminder", "ASSESSMENT")
        labels = [
            ("1", "TRAQOM digital\nattendance", BLUE),
            ("2", "Assessment digital\nattendance", TEAL),
            ("3", "Sit WA then Case\nStudy", "7C3AED"),
            ("4", "Submit answers on\nthe LMS", AMBER),
            ("5", "Sign Assessment\nSummary Record", BLUE),
        ]
        x = 0.75
        for idx, (num, label, color) in enumerate(labels):
            self.rounded(slide, x, 2.40, 2.25, 1.88, WHITE, "DCE5F0")
            self.circle(slide, x + 0.18, 2.58, 0.50, color)
            self.text(slide, x + 0.35, 2.72, 0.16, 0.18, num, 9, WHITE, True, PP_ALIGN.CENTER)
            self.text(slide, x + 0.36, 3.28, 1.54, 0.56, label, 11, DARK, True, PP_ALIGN.CENTER)
            if idx < len(labels) - 1:
                self.text(slide, x + 2.28, 3.18, 0.2, 0.4, ">", 18, GREY, True, PP_ALIGN.CENTER)
            x += 2.40
        self.text(slide, 2.4, 5.40, 8.7, 0.36, "The TRAQOM QR code and LMS submission are both handled through the LMS/TMS portal.", 13, GREY, False, PP_ALIGN.CENTER)
        self.footer(slide)

    def practice_exam_slide(self, key=None):
        slide = self.add(key)
        self.header(slide, "Take the Online AI-102 Practice Exam", "PRACTICE EXAM")
        bullets = [
            "Go to the Tertiary Exams portal and open the matching AI-102 Practice Exam.",
            "Use the practice set to test Azure AI planning, Azure OpenAI, search, vision, language and document intelligence.",
            "Review mode lets you check the answer and revisit the concept after each question.",
            "Complete the practice exam before assessment, then bring difficult items to the review session.",
        ]
        y = 2.28
        for item in bullets:
            self.text(slide, 0.92, y, 0.2, 0.25, chr(8226), 18, BLUE, True)
            self.text(slide, 1.16, y + 0.02, 4.75, 0.46, item, 13.6, DARK)
            y += 0.72
        self.text(slide, 0.92, 5.25, 4.6, 0.28, "Practice exam link:", 14, DARK, True)
        url_box = self.text(slide, 0.92, 5.58, 4.9, 0.32, EXAM_URL, 16, BLUE, True)
        try:
            url_box.text_frame.paragraphs[0].runs[0].hyperlink.address = EXAM_URL
        except Exception:
            pass
        self.rect(slide, 6.13, 2.40, 6.45, 4.70, WHITE, "CBD5E1")
        self.circle(slide, 6.30, 2.57, 0.24, BLUE)
        self.text(slide, 6.38, 2.63, 0.06, 0.08, "T", 6, WHITE, True, PP_ALIGN.CENTER)
        self.text(slide, 6.62, 2.61, 1.7, 0.18, "Tertiary Exams", 7, DARK, True)
        self.text(slide, 8.45, 2.61, 1.15, 0.18, "Practice Exams", 6, DARK)
        self.text(slide, 10.02, 2.61, 0.8, 0.18, "About Us", 6, DARK)
        self.text(slide, 11.15, 2.61, 0.5, 0.18, "Sign In", 6, DARK)
        self.rounded(slide, 11.65, 2.52, 0.58, 0.26, BLUE, BLUE)
        self.text(slide, 11.72, 2.59, 0.45, 0.08, "Get started", 5.2, WHITE, True, PP_ALIGN.CENTER)
        self.text(slide, 6.48, 3.05, 3.8, 0.18, "Portal", 5.6, BLUE, True)
        self.text(slide, 6.48, 3.23, 3.2, 0.18, "exams.tertiaryinfotech.com", 6.8, BLUE, True)
        self.text(slide, 6.48, 3.48, 3.8, 0.30, f"{COURSE_TITLE} Practice Exam", 9.0, DARK, True)
        self.text(slide, 6.48, 3.84, 3.7, 0.28, "Practice questions for Azure AI engineering, responsible AI, retrieval, agents and multimodal services.", 6.5, GREY)
        for x, title, body, color in [(6.48, "MODE", "Practice", BLUE), (7.60, "FOCUS", "AI-102", TEAL), (8.88, "ACCESS", "Portal link", "7C3AED")]:
            self.rounded(slide, x, 4.28, 1.16, 0.54, WHITE, "D6E1EF")
            self.text(slide, x + 0.16, 4.41, 0.6, 0.12, title, 5.6, GREY, True)
            self.text(slide, x + 0.16, 4.60, 0.75, 0.14, body, 6.8, DARK, True)
            self.rect(slide, x + 0.16, 4.79, 0.84, 0.03, color, color)
        self.rounded(slide, 10.75, 3.38, 1.30, 1.22, WHITE, "D6E1EF")
        self.text(slide, 10.92, 3.56, 0.82, 0.12, "PRACTICE DETAIL", 5.8, GREY, True)
        self.text(slide, 10.92, 3.82, 0.96, 0.56, "Exam code  AI-102\nUse        Review\nDomains    Azure AI\nFormat     Online", 5.3, DARK)
        self.rounded(slide, 6.48, 5.08, 1.95, 0.86, WHITE, "D6E1EF")
        self.text(slide, 6.64, 5.25, 1.25, 0.15, "Practice mode", 6.7, DARK, True)
        self.text(slide, 6.64, 5.48, 1.45, 0.24, "Check answers while revising each domain.", 5.5, GREY)
        self.rounded(slide, 8.60, 5.08, 1.95, 0.86, WHITE, "D6E1EF")
        self.text(slide, 8.76, 5.25, 1.25, 0.15, "Review before WA", 6.7, DARK, True)
        self.text(slide, 8.76, 5.48, 1.45, 0.24, "Use missed questions to guide revision.", 5.5, GREY)
        self.rounded(slide, 10.75, 5.04, 1.30, 0.92, "EAF3FF", "BBD7F2")
        self.text(slide, 10.92, 5.21, 0.88, 0.12, "Learner tip", 5.8, BLUE, True)
        self.text(slide, 10.92, 5.46, 0.90, 0.24, "Retake missed items after revision.", 5.5, GREY)
        self.footer(slide)

    def build(self) -> dict[str, int]:
        self.title_slide()
        self.section("Welcome and Housekeeping", "Course Administration", "admin_welcome")
        self.content(
            "Digital Attendance (Mandatory)",
            [
                "Take AM, PM and Assessment digital attendance for WSQ-funded courses.",
                "Scan the SSG portal QR code shown by the trainer or administrator.",
                "Minimum 75% attendance is required for funding and assessment eligibility.",
            ],
            "admin_attendance",
            "TRAQOM / SSG",
        )
        self.about_trainer("admin_about_trainer")
        self.about_trainer("admin_about_trainer_blank", blank=True)
        self.content(
            "Let's Know Each Other",
            ["Your name and organisation or role.", "Your experience with Azure or AI services.", "Your goal for this AI-102 course."],
            "admin_intro",
            "ICE-BREAKER",
        )
        self.ground_rules_slide("admin_ground_rules")
        self.two_col(
            "Lesson Plan - 2 Days",
            "Day 1",
            ["Planning, security, monitoring, responsible AI.", "Generative AI, Foundry, Azure OpenAI, RAG.", "Agentic AI and Foundry Agent Service.", "Vision, video, NLP, speech, translation."],
            "Day 2",
            ["Custom language and question answering.", "Azure AI Search, knowledge mining, vector search.", "Document Intelligence and Content Understanding.", "Capstone review and final assessment."],
            "admin_lesson_plan",
            "SCHEDULE",
        )
        self.content(
            "Skills Framework",
            [
                "TSC Title: Artificial Intelligence Application in Product Development.",
                "TSC Code: ICT-TEM-4034-1.1.",
                "Abilities: analyse algorithms, evaluate AI applications, assess feasibility and improvements.",
                "Knowledge: AI applications, performance analysis, evaluation methods, algorithm design and implementation.",
            ],
            "admin_skills_framework",
            "WSQ ALIGNMENT",
            size=14,
        )
        self.content(
            "Learning Outcomes",
            [
                "Analyse Azure AI applications and establish their correlation with efficiency.",
                "Identify and evaluate strengths and limitations of Microsoft Azure AI applications.",
                "Assess feasibility and improvements when applying Azure AI to product and maintenance processes.",
            ],
            "admin_learning_outcomes",
            "WHAT YOU WILL ACHIEVE",
        )
        self.content(
            "Course Outline",
            [f"Lab {lab.num:02d}: {lab.title}" for lab in self.labs],
            "admin_course_outline",
            "LAB-ALIGNED FLOW",
            size=13,
        )
        self.content(
            "Assessment",
            ["Written Assessment - Short-Answer Questions (WA-SAQ).", "Practical Performance (PP).", "Open book assessment: slides, Learner Guide and approved materials only.", "Final assessment starts on Day 2 at 4:00 PM."],
            "admin_assessment",
            "FINAL ASSESSMENT",
        )
        self.content(
            "Briefing for Assessment",
            ["Place phones and other materials under the table or on the floor.", "No photos or recording of assessment scripts.", "No discussion during assessment.", "Submit assessment scripts when time is up."],
            "admin_assessment_briefing",
            "ASSESSMENT RULES",
        )
        self.content(
            "Criteria for Funding",
            ["Minimum attendance rate of 75% based on SSG Digital Attendance record.", "Complete the assessment and be assessed as Competent."],
            "admin_funding",
            "SSG REQUIREMENTS",
        )
        self.content(
            "LMS / TMS and Labs",
            ["LMS/TMS: https://lms-tms.tertiaryinfotech.com", "Course labs: GitHub labs folder for TGS-2023036651.", "Use the lab sequence in this deck, Learner Guide and Lesson Plan."],
            "admin_lms_labs",
            "RESOURCES",
        )
        self.section("Day 1", "AI-102 Foundations", "day1", "Planning, responsible AI, generative AI, agents, vision and language.")
        for lab in self.labs[:6]:
            self.lab_slides(lab)
        self.section("Day 2", "AI-102 Integration and Review", "day2", "Custom language, search, document intelligence, capstone and assessment.")
        for lab in self.labs[6:]:
            self.lab_slides(lab)
        self.section("Assessment and Wrap-up", "Course Closing", "closing")
        self.practice_exam_slide("closing_practice_exam")
        self.assessment_flow_reminder("closing_assessment_flow")
        self.content("Certification and TRAQOM Survey", ["Complete the certification and TRAQOM survey to receive your certificate.", "Remember AM, PM and assessment digital attendance.", "Contact support if you need help after class."], "closing_traqom", "BEFORE YOU GO")
        self.content("Support", ["Email: enquiry@tertiaryinfotech.com", "Tel: +65 6100 0613", "Website: www.tertiarycourses.com.sg"], "closing_support", "AFTER CLASS")
        self.section("Thank You", "Congratulations", "closing_thanks", "You have designed an end-to-end Azure AI engineering solution.")
        self.prs.save(PPTX_OUT)
        shutil.copy2(PPTX_OUT, PPTX_UPDATED_OUT)
        return self.slide_index


def build_lesson_plan(labs: list[Lab], logos: list[Path], slide_index: dict[str, int]):
    doc = Document()
    setup_a4(doc)
    setup_doc_styles(doc)
    add_cover(doc, "Lesson Plan", logos)
    add_version_table(doc, "First version aligned to the generated PPT slide numbers and 10 AI-102 labs.")
    add_toc(doc, lesson_plan_toc_entries())
    doc.add_paragraph(style="Heading 1").add_run("Course Information")
    add_key_value_table(
        doc,
        [
            ("Course Title", COURSE_TITLE),
            ("Course Code", COURSE_CODE),
            ("TSC Title", TSC_TITLE),
            ("TSC Code", TSC_CODE),
            ("Duration", "2 Days (16 training hours)"),
            ("Daily Schedule", "9:00 AM - 6:00 PM, including lunch and tea breaks"),
            ("Delivery Mode", "Instructor-led, hands-on labs and guided design activities"),
            ("Assessment", "Day 2 final assessment from 4:00 PM"),
        ],
    )
    doc.add_paragraph(style="Heading 1").add_run("Lesson Schedule with Slide Numbers")
    rows = [["Time", "Topic / Activity", "Duration", "Slides"]]
    rows += [
        ["Day 1 9:00-9:20", "Digital attendance, trainer/learner introduction, ground rules", "20 min", slide_range(slide_index, ["cover", "admin_lms_labs"])],
        ["Day 1 9:20-10:20", "Lab 01 - Plan, manage, monitor, and secure an Azure AI solution", "60 min", lab_slide_range(slide_index, 1)],
        ["Day 1 10:20-10:35", "Tea break", "15 min", "-"],
        ["Day 1 10:35-11:35", "Lab 02 - Responsible AI, content safety, governance", "60 min", lab_slide_range(slide_index, 2)],
        ["Day 1 11:35-12:35", "Lab 03 - Generative AI, Foundry, Azure OpenAI, prompt flow, RAG", "60 min", lab_slide_range(slide_index, 3)],
        ["Day 1 12:35-1:00", "Review and Q&A", "25 min", lab_slide_range(slide_index, 3)],
        ["Day 1 1:00-2:00", "Lunch break", "60 min", "-"],
        ["Day 1 2:00-3:00", "Lab 04 - Agentic AI, Foundry Agent Service, multi-agent concepts", "60 min", lab_slide_range(slide_index, 4)],
        ["Day 1 3:00-4:00", "Lab 05 - Computer Vision, Custom Vision, video insights", "60 min", lab_slide_range(slide_index, 5)],
        ["Day 1 4:00-4:15", "Tea break", "15 min", "-"],
        ["Day 1 4:15-5:35", "Lab 06 - NLP, Language, Speech, Translation", "80 min", lab_slide_range(slide_index, 6)],
        ["Day 1 5:35-6:00", "Day 1 recap and checkpoint questions", "25 min", lab_slide_range(slide_index, 6)],
        ["Day 2 9:00-9:15", "Digital attendance and Day 1 recap", "15 min", str(slide_index.get("day2", "-"))],
        ["Day 2 9:15-10:15", "Lab 07 - Custom Language Models and Question Answering", "60 min", lab_slide_range(slide_index, 7)],
        ["Day 2 10:15-10:30", "Tea break", "15 min", "-"],
        ["Day 2 10:30-11:45", "Lab 08 - Azure AI Search, knowledge mining, vector search", "75 min", lab_slide_range(slide_index, 8)],
        ["Day 2 11:45-1:00", "Lab 09 - Document Intelligence and Content Understanding", "75 min", lab_slide_range(slide_index, 9)],
        ["Day 2 1:00-2:00", "Lunch break", "60 min", "-"],
        ["Day 2 2:00-3:30", "Lab 10 - AI-102 capstone and course review", "90 min", lab_slide_range(slide_index, 10)],
        ["Day 2 3:30-4:00", "Assessment briefing and review", "30 min", slide_range(slide_index, ["admin_assessment", "admin_assessment_briefing"])],
        ["Day 2 4:00-5:15", "Final assessment", "75 min", str(slide_index.get("admin_assessment", "-"))],
        ["Day 2 5:15-5:30", "Tea break", "15 min", "-"],
        ["Day 2 5:30-6:00", "TRAQOM, support, closing", "30 min", slide_range(slide_index, ["closing", "closing_thanks"])],
    ]
    add_doc_table(doc, rows)
    doc.add_paragraph(style="Heading 1").add_run("Labs Covered")
    add_doc_table(doc, [["Lab", "Main Topic", "Title", "PPT Slides"]] + [[f"{lab.num:02d}", lab.main_topic, lab.title, lab_slide_range(slide_index, lab.num)] for lab in labs])
    doc.add_paragraph(style="Heading 1").add_run("Assessment Notes")
    for item in [
        "Ensure AM, PM and Assessment digital attendance are completed.",
        "Assessment is open book using approved materials: slides and Learner Guide.",
        "Learners must complete assessment activities and be assessed as Competent for funding eligibility.",
    ]:
        render_inline_doc(doc.add_paragraph(style="List Bullet"), item)
    add_doc_footer(doc)
    doc.save(LP_DOCX_OUT)


def slide_range(slide_index: dict[str, int], keys: list[str]) -> str:
    vals = [slide_index[k] for k in keys if k in slide_index]
    if not vals:
        return "-"
    return f"{min(vals)}-{max(vals)}" if min(vals) != max(vals) else str(vals[0])


def lab_slide_range(slide_index: dict[str, int], lab_num: int) -> str:
    keys = [f"lab{lab_num:02d}_overview", f"lab{lab_num:02d}_validation"]
    return slide_range(slide_index, keys)


def main():
    OUT_DIR.mkdir(exist_ok=True)
    ASSET_DIR.mkdir(exist_ok=True)
    labs = load_labs()
    logos, _ = extract_reference_assets(labs)
    build_learner_guide(labs, logos)
    builder = DeckBuilder(labs, logos)
    slide_index = builder.build()
    build_lesson_plan(labs, logos, slide_index)
    print(f"Wrote {PPTX_OUT}")
    print(f"Wrote {LG_DOCX_OUT}")
    print(f"Wrote {LG_MD_OUT}")
    print(f"Wrote {LP_DOCX_OUT}")
    for lab in labs:
        print(f"Lab {lab.num:02d}: slides {lab_slide_range(slide_index, lab.num)}; diagram source {lab.diagram_source_slide or 'none'}")


if __name__ == "__main__":
    main()
